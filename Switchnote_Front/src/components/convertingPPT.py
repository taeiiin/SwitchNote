from flask import Flask, jsonify, request
from flask_cors import CORS
from pptx import Presentation
from urllib.parse import unquote
from dotenv import load_dotenv
import os
import boto3

app = Flask(__name__)
CORS(app)
load_dotenv()
load_dotenv('/Users/kimjuyeon/Switch_Sat/Switchnote_Front/.env')

@app.route('/convert_ppt', methods=['POST'])
def convert_ppt():
    data = request.json['data']
    category = request.json['category']
    pptTitle = request.json['title']

    # ConvertPPT 함수 실행
    urls = ConvertPPT(data, category,pptTitle)

    # URL 리스트를 JSON 형식으로 응답
    return jsonify(urls)


def ConvertPPT(data, category,pptTitle):
    #[PPT 생성]-------------------------------------------
    access_key = os.getenv('AWS_KEY')
    secret_key = os.getenv('AWS_SECRET_KEY')
    session = boto3.Session(aws_access_key_id=access_key, aws_secret_access_key=secret_key)

    # AWS S3 버킷에서 템플릿 파일 목록 가져오기
    s3 = session.client('s3')
    bucket_name = 'switchnote-ppt-template'  # S3 버킷 이름
    template_files = []
    response = s3.list_objects_v2(Bucket=bucket_name)
    for obj in response['Contents']:
        # URL 디코딩 추가
        filename = unquote(obj['Key'])
        template_files.append(filename)

    # 입력받은 category에 포함된 템플릿 파일들 가져오기
    template_paths = []
    for file in template_files:
        if category in file:
            template_paths.append(file)

    # 입력받은 category에 해당하는 템플릿 파일이 없을 경우 에러 처리
    if not template_paths:
        raise ValueError('No templates found for the given category')
    print(template_paths)
    urls = []  # 결과 URL들을 저장할 리스트
    # 템플릿 파일들 다운로드 및 PPT 생성
    for i, template_path in enumerate(template_paths):
        bucket_name = 'switchnote-ppt-template'  # ppt 템플릿 S3 버킷 이름 재언급
        # AWS S3 버킷에서 템플릿 파일 목록 가져오기
        s3 = session.client('s3')
        s3.download_file(bucket_name, template_path, f'{category}{i+1}.pptx')
        print()

        # PPT 템플릿 로드 및 PPT 생성 코드...
        presentation = Presentation(f'{category}{i+1}.pptx')
        # PPT 생성 코드...
        # 제목 페이지
        slide = presentation.slides[0]
        title = slide.shapes.title
        title.text = data[0]['title']
        subtitle_textbox = slide.placeholders[1].text_frame
        subtitle_textbox.text = data[0]['content'][0]

        # 목차 페이지
        slide = presentation.slides[1]
        content_slide = slide.placeholders[1]
        # content 요소 가져오기
        content_list = [item for item in data if item['type'] == 'b'][0]['content']
        # 텍스트로 변환하여 개행 문자 추가
        content_text = '\n'.join(content_list)
        content_slide.text = content_text

        #소제목+본문 페이지
        slide_index = 2
        previous_title = ''
        num=1
        for item in data:
            if item['type'] == 'c':
                title = item['title']
                content = item['content']

                if title != previous_title:
                    # 소제목 페이지 추가
                    slide_layout = presentation.slide_layouts[slide_index]  # 템플릿 슬라이드의 레이아웃 인덱스에 맞게 설정
                    slide = presentation.slides.add_slide(slide_layout)

                    # 소제목 슬라이드 작성
                    title_placeholder = slide.placeholders[0] #소제목
                    title_placeholder.text = title
                    number_placeholder = slide.placeholders[1] #목차번호
                    number_placeholder.text = str(num)
                    previous_title = title
                    num+=1

                # 본문 페이지 추가
                slide_layout = presentation.slide_layouts[slide_index+1]
                slide = presentation.slides.add_slide(slide_layout)


                # 소제목 작성
                title_placeholder = slide.placeholders[0] #본문 제목
                title_placeholder.text = title

                # 내용 작성
                content_placeholder = slide.placeholders[1] #본문
                content_placeholder.text = '\n'.join(content) #\n으로 각 요소 join

        # 템플릿 PPT의 마지막 5번 페이지를 새로운 PPT의 마지막에 추가
        slide = presentation.slides.add_slide(presentation.slide_layouts[-1])

        # 결과 저장
        output_path = f'{pptTitle}.pptx'
        presentation.save(output_path)

        # S3에 저장
        bucket_name = 'switchnote'  # 업로드할 S3 버킷의 이름
        s3_key = f'{pptTitle}{i+1}.pptx'  # S3 버킷에 저장될 파일의 키 (파일 경로와 이름)
        s3.upload_file(output_path, bucket_name, s3_key)

        s3_url = f"https://{bucket_name}.s3.amazonaws.com/{s3_key}"
        urls.append(s3_url)  # 생성된 PPT의 S3 URL을 리스트에 추가
        os.remove(output_path)

    return urls  # 생성된 PPT의 S3 URL 리스트 반환

if __name__ == '__main__':
    app.run(host='0.0.0.0', port=8888)