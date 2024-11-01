#!/usr/bin/env python
# coding: utf-8
#Python에서의 문서 형식 처리를 위한 aspose와 pptx api install
get_ipython().system('pip install aspose.slides')
get_ipython().system('pip install python-pptx')
#OPENAI API install
get_ipython().system('pip install openai')
from pptx import Presentation
import openai
import sys
#import spawn

def PtoS(ppt):
    #PPT 불러오기
    prs = Presentation(ppt)
    #텍스트 추출 결과 저장 리스트
    text_runs = []

    #PPT 슬라이드 순회
    for slide in prs.slides:
        #선택된 슬라이드 내에서 도형 순회
        for shape in slide.shapes:
            #도형에 텍스트 프레임이 있는 경우
            if shape.has_text_frame:
                #선택된 도형 내에서 텍스트프레임 문단 순회
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        #텍스트 저장 리스트에 저장
                        text_runs.append(run.text)
    #텍스트 데이터 정제                    
    t = []
    for i in text_runs:
        x = ' '.join(i.split())
        if x!='' or x!=' ':
            t.append(x)
    t = ' '.join(t)

    #OPENAI 연결 및 프롬프트 엔지니어링
    openai.api_key = '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'
    messages = [
        {"role": "system", "content": "You are a helpful assistant."},
        {"role": "user", "content": t},
    ]
    prompt_content = """
    주어진 텍스트를 보고 발표 대본으로 작성해 줘. 대본은 서론/본론/결론으로 나눠서 작성해주고 각 단락을 시작할 때마다 아래 형식처럼 구성하면 돼.
    <서론>
    서론 내용
    <본론>
    본론 내용
    <결론>
    결론 내용
    위와 같은 형식으로 하되, 발표자 표시 등 기타의 형식은 넣지 말고 내용이 들어갈 곳에는 발표문만 넣어줘.
    """
    messages.append({"role":"user", "content":prompt_content})

    completion_response = openai.ChatCompletion.create(
    model="gpt-3.5-turbo",
    messages=messages,
    )
    chat_res = completion_response.choices[0].message["content"]
    start = chat_res.find('<서론>')
    script = chat_res[start:]
    return(script)

#if __name__ == "__main__":
    #PtoS(sys.argv[1])